from pptx import Presentation
from pptx.util import Inches, Pt
import os

def add_slide_with_image(prs, title, content_list, image_path, img_left, img_top, img_width):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if item.startswith("  - "):
            p.text = item.replace("  - ", "")
            p.level = 1
            p.font.size = Pt(16)
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            
    slide.shapes.placeholders[1].width = Inches(5.0)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(img_left), Inches(img_top), width=Inches(img_width))
        
    return slide

def main():
    prs = Presentation()
    out_dir = "DATASET/OLAH DATA/outputs/"
    
    # SLIDE 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Strategi Segmentasi Probabilistik Calon Mahasiswa\nMenggunakan GMM & LLM"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide.placeholders[1].text = "Oleh: Mujibul Hakim (25.01.85.7010)\nMagister Teknologi Informasi, ITSNU Pekalongan"
    
    # SLIDE 2: Rumusan Masalah
    add_slide_with_image(prs, "Rumusan & Batasan Masalah", [
        "Rumusan Masalah:",
        "1. Bagaimana karakteristik segmentasi GMM berbasis IndoBERT 2019–2024?",
        "2. Bagaimana evolusi stabilitas klaster (ARI, Jaccard) akibat pandemi COVID-19?",
        "3. Bagaimana otomasi LLM (OpenRouter/Llama 3.3) mengoptimasi generasi persona?",
        "4. Bagaimana validitas otomasi LLM berdasarkan penilaian pakar (Expert Judgement)?",
        "Batasan:",
        "Menggunakan data sekunder 2019-2024 (nama, asal sekolah, alamat) dan Hybrid Cognitive Pipeline LLM."
    ], os.path.join(out_dir, "gambar_4_1_distribusi.png"), 5.5, 2.0, 4.0)
    
    # SLIDE 3: Evaluasi Klaster K
    add_slide_with_image(prs, "Penentuan Klaster Optimal (GMM)", [
        "Hasil Integrasi IndoBERT & GMM (2019-2024):",
        "  - K optimal bervariasi tiap periode (6, 6, 6, 5, 2, 3) klaster.",
        "  - Terjadi konsolidasi segmen pendaftar yang homogen dari 6 klaster (2019) menjadi 2 klaster (2023).",
        "  - Silhouette Score (misal 0.0905 pada 2023) menunjukkan model probabilitas GMM efektif menangani tumpang tindih segmen dibanding hard clustering."
    ], os.path.join(out_dir, "gambar_4_3a_silhouette.png"), 5.0, 2.0, 4.5)
    
    # SLIDE 4: Evaluasi Time Series
    add_slide_with_image(prs, "Stabilitas Model & Structural Breaks", [
        "Dinamika Perubahan Struktural Pendaftar:",
        "  - Evaluasi stabilitas (ARI) membuktikan asumsi pasar statis adalah keliru.",
        "  - Seluruh 5 transisi antar periode (2019-2024) adalah structural break (ARI < 0.30).",
        "  - Transisi Baseline ke COVID Crisis (2019-2020) mencatat disrupsi ekstrem (ARI = -0.0036).",
        "  - Model GMM sukses memetakan pergeseran centroid (Centroid Drift) ini secara presisi."
    ], os.path.join(out_dir, "gambar_4_3c_ari.png"), 5.0, 2.0, 4.5)
    
    # SLIDE 5: Generasi Persona LLM
    add_slide_with_image(prs, "Validasi Persona LLM (Hybrid Pipeline)", [
        "Otomasi Penalaran (Reasoning) Skala Besar:",
        "  - Menggunakan OpenRouter (Llama 3.3 70B, Nemotron) mengubah data kuantitatif ke narasi.",
        "  - LLM berhasil men-generate narasi persona calon mahasiswa dengan logika kausal yang tajam.",
        "  - Divalidasi menggunakan Expert Judgement, membuktikan kapabilitas Generative AI untuk ranah manajerial ITSNU."
    ], os.path.join(out_dir, "gambar_4_5_proyeksi.png"), 5.0, 2.0, 4.5)
    
    # SLIDE 6: Kesimpulan & Saran
    add_slide_with_image(prs, "Kesimpulan & Rekomendasi", [
        "Kesimpulan Utama:",
        "1. IndoBERT + GMM sangat efektif memetakan probabilitas pendaftar tanpa paksaan satu kategori absolut.",
        "2. Pasar PMB ITSNU sangat dinamis (ARI < 0.30), sehingga marketing butuh strategi probabilistik yang adaptif.",
        "3. Hybrid Cognitive Pipeline (LLM) tervalidasi akurat mentransformasi angka ke narasi manajerial marketing.",
        "Rekomendasi:",
        "Terapkan dashboard PMB-LLM untuk merancang strategi marketing real-time berdasarkan evolusi pendaftar."
    ], "", 5.0, 2.0, 4.5)
    
    prs.save("FULL TESIS/Presentasi_Tesis_WS2_FINAL_SINKRON.pptx")
    print("Presentasi_Tesis_WS2_FINAL_SINKRON.pptx created!")

if __name__ == '__main__':
    main()
