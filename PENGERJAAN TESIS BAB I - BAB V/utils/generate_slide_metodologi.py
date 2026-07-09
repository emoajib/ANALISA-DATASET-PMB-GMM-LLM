from pptx import Presentation
from pptx.util import Inches, Pt

def main():
    prs = Presentation()
    
    # Tambahkan 1 slide untuk Metodologi
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Judul
    slide.shapes.title.text = "Metodologi Penelitian (CRISP-DM & Pipeline)"
    
    # Isi Teks
    tf = slide.shapes.placeholders[1].text_frame
    
    bullets = [
        "Pendekatan Standar: Mengadopsi kerangka kerja Cross-Industry Standard Process for Data Mining (CRISP-DM).",
        "Arsitektur Hybrid Cognitive Pipeline:",
        "  - 1. Pre-processing & Embedding: Mengubah data teks (Nama, Sekolah, Alamat) menjadi vektor 768D (IndoBERT).",
        "  - 2. Reduksi Dimensi: Menggunakan PCA untuk mengoptimasi dimensi vektor.",
        "  - 3. Klasterisasi Probabilistik: Gaussian Mixture Model (GMM) untuk menghitung probabilitas posterior tiap mahasiswa.",
        "  - 4. Generasi Persona: LLM skala besar (Llama 3.3 via OpenRouter) mengekstrak narasi manajerial dari centroid GMM."
    ]
    
    for i, item in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if item.startswith("  - "):
            p.text = item.replace("  - ", "")
            p.level = 1
            p.font.size = Pt(18)
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(20)
            
    # Tambahkan Notes (Narasi Video)
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    text_frame.text = "Untuk menjawab rumusan masalah tersebut, penelitian ini mengadopsi standar metodologi CRISP-DM dengan arsitektur Hybrid. Data teks pendaftar pertama-tama diubah menjadi angka 768 dimensi oleh IndoBERT, lalu direduksi menggunakan PCA, dan diklasterisasi secara probabilistik menggunakan Gaussian Mixture Model. Hasil matematis dari GMM ini kemudian diproses oleh LLM OpenRouter untuk menghasilkan persona bahasa manusia."
    
    prs.save("FULL TESIS/Tambahan_Slide_Metodologi.pptx")
    print("Tambahan_Slide_Metodologi.pptx created successfully!")

if __name__ == '__main__':
    main()
