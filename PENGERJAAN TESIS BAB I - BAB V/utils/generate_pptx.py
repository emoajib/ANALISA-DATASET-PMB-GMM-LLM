from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def add_slide(prs, title, content_list):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    
    body_shape = slide.shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(22)
        else:
            p = tf.add_paragraph()
            if item.startswith("  - "):
                p.text = item.replace("  - ", "")
                p.level = 1
                p.font.size = Pt(20)
            else:
                p.text = item
                p.level = 0
                p.font.size = Pt(22)
    return slide

def main():
    prs = Presentation()
    
    # SLIDE 1: Title
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Strategi Segmentasi Probabilistik Calon Mahasiswa\nMenggunakan GMM & LLM"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    subtitle.text = "Oleh: Mujibul Hakim (25.01.85.7010)\nMagister Teknologi Informasi, ITSNU Pekalongan"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    
    # SLIDE 2: Latar Belakang
    add_slide(prs, "Latar Belakang Masalah", [
        "Data penerimaan mahasiswa baru (PMB) ITSNU Pekalongan belum dimanfaatkan secara optimal untuk merumuskan strategi pemasaran.",
        "Segmentasi saat ini masih bersifat deskriptif dasar, belum menggunakan pemodelan Machine Learning.",
        "Dibutuhkan pendekatan Hybrid yang menggabungkan:",
        "  - Natural Language Processing (IndoBERT) untuk data teks",
        "  - Gaussian Mixture Model (GMM) untuk klasterisasi probabilistik",
        "  - Large Language Model (LLM) untuk mendefinisikan persona mahasiswa"
    ])
    
    # SLIDE 3: Rumusan Masalah
    add_slide(prs, "Rumusan & Batasan Masalah", [
        "Rumusan Masalah:",
        "  - Bagaimana mengintegrasikan model bahasa IndoBERT dengan algoritma GMM untuk segmentasi profil mahasiswa?",
        "  - Bagaimana GMM menangani variabilitas demografis berbasis Time-Series (2019-2024)?",
        "  - Bagaimana performa LLM dalam mengotomasi pelabelan profil (persona) klaster?",
        "Batasan Masalah:",
        "  - Menggunakan data sekunder PMB ITSNU tahun 2019 hingga 2024.",
        "  - Fokus pada fitur tekstual (Nama, Alamat, Asal Sekolah, Kab/Kota)."
    ])
    
    # SLIDE 4: Metodologi (CRISP-DM)
    add_slide(prs, "Metodologi Penelitian: CRISP-DM", [
        "Penelitian ini mengadopsi standar Cross-Industry Standard Process for Data Mining (CRISP-DM):",
        "  1. Business Understanding: Identifikasi kebutuhan PMB",
        "  2. Data Understanding: Eksplorasi atribut data (2019-2024)",
        "  3. Data Preparation: Cleansing & Ekstraksi Teks",
        "  4. Modeling: IndoBERT Embedding & GMM Clustering",
        "  5. Evaluation: K-Scan (BIC/AIC/Silhouette) & ARI Time-Series",
        "  6. Deployment: Pembuatan Dashboard Streamlit"
    ])
    
    # SLIDE 5: Arsitektur Sistem
    add_slide(prs, "Arsitektur Sistem (Pipeline)", [
        "Sistem dibangun dengan arsitektur Hybrid:",
        "1. Ekstraksi Fitur Teks: IndoBERT-Base-P1 mengubah data teks menjadi vektor 768 dimensi.",
        "2. Reduksi Dimensi: PCA mereduksi 768D menjadi komponen utama untuk optimasi.",
        "3. Klasterisasi: Gaussian Mixture Model (GMM) membentuk segmentasi berdasarkan probabilitas.",
        "4. Generasi Persona: Prompt Engineering pada Ollama (LLM lokal) untuk membaca deskripsi statistik klaster dan menghasilkan 'Persona' target audiens."
    ])
    
    # SLIDE 6: Evaluasi Penentuan Klaster (K)
    add_slide(prs, "Evaluasi Klaster (K-Scan)", [
        "Penentuan jumlah klaster optimal (K) dievaluasi menggunakan:",
        "  - Bayesian Information Criterion (BIC)",
        "  - Akaike Information Criterion (AIC)",
        "  - Silhouette Coefficient",
        "Hasil:",
        "Grafik BIC dan AIC menunjukkan titik elbow/penurunan drastis pada rentang K=3 hingga K=5. Evaluasi kestabilan menunjukkan K yang optimal adalah 4 klaster utama."
    ])
    
    # SLIDE 7: Profil Klaster (Otomasi LLM)
    add_slide(prs, "Profil & Persona Klaster (LLM)", [
        "LLM mengidentifikasi 4 persona utama calon mahasiswa:",
        "  - Klaster 0: Mahasiswa Teknis Lokal (Dominan prodi Teknologi Informasi dari SMK lokal).",
        "  - Klaster 1: Pencari Stabilitas Karir (Dominan prodi K3 dari SMA/SMK dengan ketertarikan industri).",
        "  - Klaster 2: Penggerak Ekonomi Desa (Fokus prodi Agribisnis dari MA/SMA pinggiran).",
        "  - Klaster 3: Inovator Lintas Disiplin (Peminat prodi Teknik Industri dari wilayah urban).",
        "[Tempatkan Tabel 4.18 Perbandingan LLM di sini]"
    ])
    
    # SLIDE 8: Stabilitas Time-Series
    add_slide(prs, "Stabilitas Model (Time-Series 2019-2024)", [
        "Evaluasi konsistensi profil lintas waktu menggunakan Adjusted Rand Index (ARI):",
        "  - Model menunjukkan stabilitas profil dari tahun ke tahun meskipun terdapat fluktuasi jumlah pendaftar.",
        "  - Terjadi pergeseran centroid (Centroid Drift) minor pada tahun 2021-2022 akibat efek pandemi.",
        "  - Secara keseluruhan, GMM mampu beradaptasi dengan distribusi baru tanpa kehilangan bentuk klaster utama."
    ])
    
    # SLIDE 9: Kesimpulan
    add_slide(prs, "Kesimpulan Utama", [
        "1. Integrasi IndoBERT dan GMM berhasil mengidentifikasi 4 klaster tersembunyi dengan skor evaluasi yang solid (Silhouette > 0.4 pada reduksi PCA).",
        "2. Otomasi pelabelan menggunakan LLM lokal (Ollama) terbukti akurat dalam menghasilkan persona marketing yang deskriptif dan logis.",
        "3. Evaluasi Time-Series menunjukkan strategi segmentasi ini sangat andal (robust) untuk digunakan memproyeksi pendaftar di masa depan."
    ])
    
    # SLIDE 10: Saran & Rekomendasi 2025
    add_slide(prs, "Saran & Rekomendasi Marketing", [
        "Rekomendasi Taktis PMB:",
        "  - Klaster Lokal (TI): Gunakan media sosial (Instagram) dengan konten teknologi & sertifikasi.",
        "  - Klaster K3: Bekerjasama dengan kawasan industri untuk sosialisasi K3.",
        "  - Klaster Agribisnis: Pendekatan langsung ke desa/sekolah pinggiran tentang digitalisasi tani.",
        "Saran Penelitian Lanjutan:",
        "  - Mengintegrasikan data akademik (nilai rapor) dan sosial ekonomi.",
        "  - Membandingkan GMM dengan arsitektur Deep Clustering (misal: DEC)."
    ])
    
    # SLIDE 11: Penutup
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Terima Kasih"
    subtitle.text = "Sesi Tanya Jawab (Q&A)"
    
    prs.save("FULL TESIS/Presentasi_Tesis_WS2.pptx")
    print("PPTX generated!")

if __name__ == '__main__':
    main()
