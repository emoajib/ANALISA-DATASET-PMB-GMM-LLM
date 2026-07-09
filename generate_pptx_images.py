from pptx import Presentation
from pptx.util import Inches, Pt
import os
import csv

def add_slide_with_image(prs, title, content_list, image_path, img_left, img_top, img_width):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    
    for i, item in enumerate(content_list):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if item.startswith("  - "):
            p.text = item.replace("  - ", "")
            p.level = 1
            p.font.size = Pt(16)
        else:
            p.text = item
            p.level = 0
            p.font.size = Pt(18)
            
    # Resize text box to make room for image
    slide.shapes.placeholders[1].width = Inches(5.0)
    
    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(img_left), Inches(img_top), width=Inches(img_width))
        
    return slide

def add_slide_with_table(prs, title, csv_path):
    slide_layout = prs.slide_layouts[5] # Title only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    
    if os.path.exists(csv_path):
        with open(csv_path, 'r', encoding='utf-8') as f:
            reader = list(csv.reader(f))
            if not reader: return
            
            rows = len(reader)
            cols = len(reader[0])
            
            # Reduce rows if too large for slide (max 6 rows for presentation)
            display_rows = reader[:7]
            
            x, y, cx, cy = Inches(0.5), Inches(2.0), Inches(9.0), Inches(4.5)
            shape = slide.shapes.add_table(len(display_rows), cols, x, y, cx, cy)
            table = shape.table
            
            for i, row in enumerate(display_rows):
                for j, cell_val in enumerate(row):
                    cell = table.cell(i, j)
                    cell.text = cell_val
                    cell.text_frame.paragraphs[0].font.size = Pt(12)

def main():
    prs = Presentation()
    out_dir = "DATASET/OLAH DATA/outputs/"
    
    # SLIDE 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Strategi Segmentasi Probabilistik Calon Mahasiswa\nMenggunakan GMM & LLM"
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    slide.placeholders[1].text = "Oleh: Mujibul Hakim (25.01.85.7010)\nMagister Teknologi Informasi, ITSNU Pekalongan"
    
    # SLIDE 2
    add_slide_with_image(prs, "Latar Belakang Masalah", [
        "Distribusi demografis PMB ITSNU (kanan) butuh strategi pemetaan otomatis.",
        "Segmentasi saat ini masih bersifat deskriptif dasar.",
        "Pendekatan Hybrid:",
        "  - IndoBERT (Teks)",
        "  - GMM (Klasterisasi)",
        "  - LLM (Persona)"
    ], os.path.join(out_dir, "gambar_4_1_distribusi.png"), 5.5, 2.0, 4.0)
    
    # SLIDE 3: Evaluasi Klaster
    add_slide_with_image(prs, "Evaluasi Klaster (K-Scan)", [
        "Metode penentuan klaster:",
        "  - BIC & AIC",
        "  - Silhouette Coefficient",
        "Grafik Silhouette Score (kanan) menunjukkan K=4 sebagai jumlah klaster optimal dengan pemisahan spasial terbaik."
    ], os.path.join(out_dir, "gambar_4_3a_silhouette.png"), 5.0, 2.0, 4.5)
    
    # SLIDE 4: Scatter Plot GMM 2024
    add_slide_with_image(prs, "Distribusi Klaster GMM (2024)", [
        "Visualisasi Scatter Plot (PCA 2D) untuk pendaftar tahun 2024:",
        "  - Klaster 0: Mahasiswa Teknis Lokal",
        "  - Klaster 1: Pencari Stabilitas Karir",
        "  - Klaster 2: Penggerak Ekonomi Desa",
        "  - Klaster 3: Inovator Lintas Disiplin"
    ], os.path.join(out_dir, "gambar_4_2f_scatter_2024.png"), 5.0, 2.0, 4.5)
    
    # SLIDE 5: Evaluasi Time Series
    add_slide_with_image(prs, "Stabilitas Time-Series (2019-2024)", [
        "Evaluasi konsistensi antar tahun (Adjusted Rand Index):",
        "  - Nilai ARI konstan stabil di atas 0.70",
        "  - Menandakan GMM kebal terhadap fluktuasi minor pendaftar (robust).",
        "  - Pergeseran centroid minor terjadi akibat efek pandemi, namun tidak merusak profil utama."
    ], os.path.join(out_dir, "gambar_4_3c_ari.png"), 5.0, 2.0, 4.5)
    
    # SLIDE 6: Tabel Persona LLM
    add_slide_with_table(prs, "Otomasi Profiling LLM (Persona)", os.path.join(out_dir, "tabel_4_18_perbandingan.csv"))
    
    # SLIDE 7: Proyeksi & Kesimpulan
    add_slide_with_image(prs, "Proyeksi & Kesimpulan", [
        "Kesimpulan:",
        "1. IndoBERT & GMM akurat membentuk 4 klaster utama.",
        "2. LLM lokal (Ollama) sangat efektif merumuskan persona target.",
        "3. Sistem stabil lintas waktu (ARI tinggi).",
        "Saran (Proyeksi 2025):",
        "Gunakan strategi promosi yang spesifik per klaster seperti pada grafik (kanan)."
    ], os.path.join(out_dir, "gambar_4_5_proyeksi.png"), 5.0, 2.0, 4.5)
    
    prs.save("FULL TESIS/Presentasi_Tesis_WS2_Bergambar.pptx")
    print("Presentasi_Tesis_WS2_Bergambar.pptx created!")

if __name__ == '__main__':
    main()
