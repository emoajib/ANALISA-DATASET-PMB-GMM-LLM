from pptx import Presentation

def add_notes():
    prs = Presentation("FULL TESIS/Presentasi_Tesis_WS2_FINAL_SINKRON.pptx")
    
    naras = [
        # Slide 1: Cover
        "Assalamualaikum Warahmatullahi Wabarakatuh. Yang terhormat Tim Reviewer dan Dosen Pembimbing. Perkenalkan saya Mujibul Hakim. Pada kesempatan ini saya akan mempresentasikan tesis saya yang berjudul 'Strategi Segmentasi Probabilistik Calon Mahasiswa Menggunakan GMM dan LLM'.",
        
        # Slide 2: Rumusan & Batasan Masalah
        "Latar belakang penelitian ini bermula dari belum optimalnya data PMB di ITSNU Pekalongan. Oleh karena itu, rumusan masalah berfokus pada 4 hal: Pertama, pembentukan klaster GMM berbasis IndoBERT. Kedua, evaluasi stabilitas klaster saat pandemi. Ketiga, otomasi generasi persona dengan LLM. Dan keempat, validasi pakar. Batasan masalah hanya menggunakan data sekunder dari 2019 hingga 2024.",
        
        # Slide 3: Penentuan Klaster
        "Masuk ke hasil pertama. Integrasi teks IndoBERT dengan GMM berhasil membentuk segmentasi probabilitas yang dinamis. Dari analisis K-Scan menggunakan grafik Silhouette, kita bisa melihat bahwa jumlah klaster optimal bervariasi setiap tahunnya, dan perlahan mengerucut menjadi 2 klaster utama di tahun 2023. GMM terbukti sangat efektif mengukur probabilitas tumpang tindih segmen ini.",
        
        # Slide 4: Evaluasi Time Series
        "Namun, hasil paling mengejutkan terlihat pada analisis stabilitas Time-Series. Menggunakan Adjusted Rand Index (ARI), terbukti bahwa asumsi pasar pendidikan itu statis adalah keliru. Terjadi pergeseran ekstrem atau structural break pada transisi tahun 2019 ke 2020 akibat pandemi COVID-19 dengan nilai ARI negatif. GMM berhasil menangkap pergeseran centroid ini secara akurat.",
        
        # Slide 5: Generasi Persona LLM
        "Setelah klaster terbentuk, sistem LLM Hybrid Cognitive Pipeline, khususnya model OpenRouter Llama 3.3, digunakan untuk menerjemahkan angka kuantitatif tersebut menjadi narasi persona. LLM mampu mendeskripsikan siapa sebenarnya audiens kita di setiap klaster, dan hasil narasi ini telah dinyatakan sangat valid berdasarkan penilaian Expert Judgement atau pakar.",
        
        # Slide 6: Kesimpulan
        "Sebagai penutup, tiga kesimpulan utama penelitian ini adalah: Pertama, GMM dan IndoBERT akurat memetakan profil tanpa paksaan kategori mutlak. Kedua, pasar PMB sangat rentan terhadap perubahan dinamis sehingga butuh strategi adaptif. Ketiga, LLM sangat valid merumuskan narasi manajerial. Ke depannya, pihak kampus dapat memanfaatkan dashboard cerdas ini untuk merancang kampanye marketing secara presisi per klaster."
    ]
    
    for i, slide in enumerate(prs.slides):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = naras[i] if i < len(naras) else ""
        
    prs.save("FULL TESIS/Presentasi_Tesis_WS2_Dengan_Narasi.pptx")
    print("Narasi added successfully to Presentasi_Tesis_WS2_Dengan_Narasi.pptx")

if __name__ == '__main__':
    add_notes()
